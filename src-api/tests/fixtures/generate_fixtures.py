"""Generate binary test fixture files for document parser tests.

Run once from src-api/:
    uv run python tests/fixtures/generate_fixtures.py
"""
from pathlib import Path

HERE = Path(__file__).parent

NAME = "John Doe"
TITLE = "Senior Software Engineer"
SUMMARY = "Experienced professional with expertise in Python, FastAPI, and cloud infrastructure."
SKILLS = ["Python", "FastAPI", "PostgreSQL", "Docker"]
COMPANY = "Acme Corp"


def generate_docx():
    from docx import Document

    doc = Document()
    doc.add_heading(NAME, level=1)
    doc.add_heading(TITLE, level=2)
    doc.add_paragraph(SUMMARY)
    doc.add_heading("Skills", level=3)
    for skill in SKILLS:
        doc.add_paragraph(skill, style="List Bullet")
    doc.add_heading("Experience", level=3)
    doc.add_paragraph(f"{COMPANY} — {TITLE} (2020-2024)")
    doc.save(HERE / "sample.docx")
    print("Generated sample.docx")


def generate_pdf():
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    text = f"{NAME}\n{TITLE}\n\n{SUMMARY}\n\nSkills: {', '.join(SKILLS)}\n\n{COMPANY} — {TITLE} (2020-2024)"
    page.insert_text((72, 72), text, fontsize=12)
    doc.save(str(HERE / "sample.pdf"))
    doc.close()
    print("Generated sample.pdf")


def generate_xlsx():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Resume"
    ws["A1"] = "Name"
    ws["B1"] = NAME
    ws["A2"] = "Title"
    ws["B2"] = TITLE
    ws["A3"] = "Summary"
    ws["B3"] = SUMMARY
    ws["A4"] = "Skills"
    ws["B4"] = ", ".join(SKILLS)
    ws["A5"] = "Company"
    ws["B5"] = COMPANY
    wb.save(HERE / "sample.xlsx")
    print("Generated sample.xlsx")


def generate_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.text = f"{NAME}\n{TITLE}\n{SUMMARY}\nSkills: {', '.join(SKILLS)}\n{COMPANY}"
    prs.save(HERE / "sample.pptx")
    print("Generated sample.pptx")


if __name__ == "__main__":
    generate_docx()
    generate_pdf()
    generate_xlsx()
    generate_pptx()
    print("All fixtures generated.")

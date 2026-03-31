from pathlib import Path

import markdown
from bs4 import BeautifulSoup

ALLOWED_CONTENT_TYPES = {
    "text/markdown": ".md",
    "text/x-markdown": ".md",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}

EXTENSION_TO_CONTENT_TYPE = {v: k for k, v in ALLOWED_CONTENT_TYPES.items() if k != "text/x-markdown"}


def resolve_content_type(filename: str, declared_content_type: str) -> str:
    """Resolve content type from filename extension if declared type is generic."""
    if declared_content_type not in ("application/octet-stream", ""):
        return declared_content_type
    ext = Path(filename).suffix.lower()
    return EXTENSION_TO_CONTENT_TYPE.get(ext, declared_content_type)


def parse_document(file_path: str, content_type: str) -> str:
    """Parse a document and return extracted plain text."""
    parsers = {
        "text/markdown": _parse_markdown,
        "text/x-markdown": _parse_markdown,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _parse_docx,
        "application/pdf": _parse_pdf,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _parse_xlsx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": _parse_pptx,
    }
    parser = parsers.get(content_type)
    if parser is None:
        raise ValueError(f"Unsupported content type: {content_type}")
    return parser(file_path)


def _parse_markdown(file_path: str) -> str:
    text = Path(file_path).read_text(encoding="utf-8")
    html = markdown.markdown(text)
    return BeautifulSoup(html, "html.parser").get_text(separator="\n").strip()


def _parse_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _parse_pdf(file_path: str) -> str:
    import fitz

    doc = fitz.open(file_path)
    parts = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            parts.append(text)
    doc.close()
    return "\n\n".join(parts)


def _parse_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())
    return "\n".join(parts)
